"""
chunks.py
---------
Document extraction and chunking layer for the AI knowledge assistant.

Responsibilities:
- Read uploaded knowledge-source files.
- Extract text from supported document formats.
- Split extracted text into overlapping chunks.
- Store chunks in MongoDB.
- Keep chunks isolated between users.

Pipeline:

    upload.py
        ↓
    knowledge_sources
        ↓
    chunks.py
        ↓
    knowledge_chunks
        ↓
    embedding / vector_index.py

Important:
- This file does NOT create embeddings.
- This file does NOT perform vector search.
- This file does NOT create the Knowledge Graph.
"""

from __future__ import annotations

import csv
import re
from pathlib import Path
from typing import Any

from bs4 import BeautifulSoup
from docx import Document
from fastapi import HTTPException
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

from config import settings
from database import (
    create_knowledge_chunks,
    get_knowledge_source,
    knowledge_chunks,
)
from embeddings import get_embeddings_batch


# ============================================================
# Text Extraction
# ============================================================

def extract_text_from_pdf(
    file_path: Path,
) -> str:
    """Extract text from a PDF file."""

    reader = PdfReader(
        str(file_path)
    )

    pages: list[str] = []

    for page in reader.pages:
        text = page.extract_text()

        if text:
            pages.append(text)

    return "\n\n".join(pages).strip()


def extract_text_from_docx(
    file_path: Path,
) -> str:
    """Extract text from a DOCX file."""

    document = Document(
        str(file_path)
    )

    paragraphs: list[str] = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()

        if text:
            paragraphs.append(text)

    # --------------------------------------------------------
    # Extract basic table content as well.
    # --------------------------------------------------------

    for table in document.tables:
        for row in table.rows:
            cells = [
                cell.text.strip()
                for cell in row.cells
            ]

            row_text = " | ".join(
                cell
                for cell in cells
                if cell
            )

            if row_text:
                paragraphs.append(
                    row_text
                )

    return "\n".join(paragraphs).strip()


def extract_text_from_xlsx(
    file_path: Path,
) -> str:
    """Extract readable text from an XLSX workbook."""

    workbook = load_workbook(
        filename=str(file_path),
        read_only=True,
        data_only=True,
    )

    sections: list[str] = []

    try:
        for worksheet in workbook.worksheets:
            sections.append(
                f"=== Sheet: {worksheet.title} ==="
            )

            for row in worksheet.iter_rows(
                values_only=True
            ):
                values = []

                for value in row:
                    if value is None:
                        values.append("")
                    else:
                        values.append(
                            str(value).strip()
                        )

                row_text = " | ".join(
                    values
                ).strip()

                if row_text:
                    sections.append(
                        row_text
                    )
    finally:
        workbook.close()

    return "\n".join(
        sections
    ).strip()


def extract_text_from_xls(
    file_path: Path,
) -> str:
    """
    Extract text from legacy XLS files.

    XLS support requires an appropriate Excel reader.
    This function intentionally raises a clear error rather
    than silently producing incorrect data.
    """

    raise HTTPException(
        status_code=400,
        detail=(
            "Legacy .xls files are not supported yet. "
            "Please convert the file to .xlsx and upload it again."
        ),
    )


def extract_text_from_pptx(
    file_path: Path,
) -> str:
    """Extract text from a PowerPoint presentation."""

    presentation = Presentation(
        str(file_path)
    )

    slides: list[str] = []

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1,
    ):
        slide_parts: list[str] = [
            f"=== Slide {slide_number} ==="
        ]

        for shape in slide.shapes:
            if not hasattr(
                shape,
                "text",
            ):
                continue

            text = shape.text.strip()

            if text:
                slide_parts.append(
                    text
                )

        slides.append(
            "\n".join(slide_parts)
        )

    return "\n\n".join(
        slides
    ).strip()


def extract_text_from_txt(
    file_path: Path,
) -> str:
    """Extract text from plain-text files."""

    return file_path.read_text(
        encoding="utf-8",
        errors="replace",
    ).strip()


def extract_text_from_csv(
    file_path: Path,
) -> str:
    """Extract readable text from a CSV file."""

    sections: list[str] = []

    with file_path.open(
        "r",
        encoding="utf-8",
        errors="replace",
        newline="",
    ) as csv_file:

        reader = csv.reader(
            csv_file
        )

        for row in reader:
            row_text = " | ".join(
                cell.strip()
                for cell in row
            ).strip()

            if row_text:
                sections.append(
                    row_text
                )

    return "\n".join(
        sections
    ).strip()


def extract_text_from_html(
    file_path: Path,
) -> str:
    """Extract visible text from HTML."""

    html = file_path.read_text(
        encoding="utf-8",
        errors="replace",
    )

    soup = BeautifulSoup(
        html,
        "lxml",
    )

    for element in soup(
        [
            "script",
            "style",
            "noscript",
        ]
    ):
        element.decompose()

    return soup.get_text(
        separator="\n",
        strip=True,
    )


def extract_text_from_xml(
    file_path: Path,
) -> str:
    """Extract readable text from XML."""

    xml = file_path.read_text(
        encoding="utf-8",
        errors="replace",
    )

    soup = BeautifulSoup(
        xml,
        "xml",
    )

    return soup.get_text(
        separator="\n",
        strip=True,
    )


# ============================================================
# Generic Document Extraction
# ============================================================

def extract_text(
    file_path: str,
    file_type: str,
) -> str:
    """
    Extract text from a supported document.

    Args:
        file_path:
            Path to the uploaded file.

        file_type:
            Normalized file type stored by upload.py.

    Returns:
        Extracted text.
    """

    path = Path(
        file_path
    )

    if not path.exists():
        raise HTTPException(
            status_code=404,
            detail="Uploaded file could not be found.",
        )

    extractors = {
        "pdf": extract_text_from_pdf,
        "docx": extract_text_from_docx,
        "xlsx": extract_text_from_xlsx,
        "xls": extract_text_from_xls,
        "pptx": extract_text_from_pptx,
        "txt": extract_text_from_txt,
        "markdown": extract_text_from_txt,
        "csv": extract_text_from_csv,
        "html": extract_text_from_html,
        "xml": extract_text_from_xml,
    }

    extractor = extractors.get(
        file_type
    )

    if extractor is None:
        raise HTTPException(
            status_code=400,
            detail=(
                f"No text extractor is available "
                f"for file type '{file_type}'."
            ),
        )

    try:
        text = extractor(
            path
        )
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(
            status_code=422,
            detail=(
                f"Failed to extract text from "
                f"the uploaded file: {exc}"
            ),
        ) from exc

    text = text.strip()

    if not text:
        raise HTTPException(
            status_code=422,
            detail=(
                "No readable text was found in the uploaded file."
            ),
        )

    return text


# ============================================================
# Text Normalization
# ============================================================

def normalize_text(
    text: str,
) -> str:
    """
    Normalize extracted document text.

    This removes excessive whitespace while preserving
    paragraph boundaries.
    """

    lines = [
        " ".join(
            line.split()
        )
        for line in text.splitlines()
    ]

    cleaned_lines = [
        line
        for line in lines
        if line
    ]

    return "\n".join(
        cleaned_lines
    ).strip()


# ============================================================
# Text Chunking
# ============================================================

def chunk_text(
    text: str,
    chunk_size: int | None = None,
    chunk_overlap: int | None = None,
) -> list[str]:
    """
    Split text into overlapping chunks at sentence boundaries.

    Never breaks mid-sentence or mid-word. Respects paragraph
    boundaries and sentence endings.
    """

    if chunk_size is None:
        chunk_size = settings.CHUNK_SIZE

    if chunk_overlap is None:
        chunk_overlap = settings.CHUNK_OVERLAP

    if chunk_size <= 0:
        raise ValueError("chunk_size must be greater than zero.")
    if chunk_overlap < 0:
        raise ValueError("chunk_overlap cannot be negative.")
    if chunk_overlap >= chunk_size:
        raise ValueError("chunk_overlap must be smaller than chunk_size.")

    text = normalize_text(text)
    if not text:
        return []

    # Split into sentences first
    sentences = re.split(r'(?<=[.!?])\s+', text.strip())
    sentences = [s.strip() for s in sentences if s.strip()]

    if not sentences:
        return []

    chunks: list[str] = []
    current_chunk: list[str] = []
    current_length = 0

    for sentence in sentences:
        sentence_length = len(sentence)

        # If adding this sentence would exceed chunk_size, save current chunk
        if current_length + sentence_length + 1 > chunk_size and current_chunk:
            # Save current chunk
            chunk_text_str = ' '.join(current_chunk).strip()
            if chunk_text_str:
                chunks.append(chunk_text_str)

            # Start new chunk with overlap
            # Find sentences to overlap (from end of current chunk)
            overlap_text = ' '.join(current_chunk)
            overlap_start = max(0, len(overlap_text) - chunk_overlap)

            # Find the last sentence that fits in overlap
            overlap_sentences = []
            overlap_length = 0
            for sent in reversed(current_chunk):
                if overlap_length + len(sent) + 1 <= chunk_overlap:
                    overlap_sentences.insert(0, sent)
                    overlap_length += len(sent) + 1
                else:
                    break

            current_chunk = overlap_sentences
            current_length = sum(len(s) for s in current_chunk) + len(current_chunk) - 1

        # Add sentence to current chunk
        current_chunk.append(sentence)
        current_length += sentence_length + 1

    # Add final chunk
    if current_chunk:
        chunk_text_str = ' '.join(current_chunk).strip()
        if chunk_text_str:
            chunks.append(chunk_text_str)

    return chunks


# ============================================================
# Create Chunks For Source
# ============================================================

def process_knowledge_source(
    source_id: str,
    user_id: str,
) -> list[str]:
    """
    Extract and chunk a user's uploaded document.

    The source must belong to the requesting user.

    Returns:
        List of created MongoDB chunk IDs.
    """

    source = get_knowledge_source(
        source_id=source_id,
        user_id=user_id,
    )

    if source is None:
        raise HTTPException(
            status_code=404,
            detail="Knowledge source not found.",
        )

    # --------------------------------------------------------
    # Extract document text
    # --------------------------------------------------------

    text = extract_text(
        file_path=source["file_path"],
        file_type=source["file_type"],
    )

    # --------------------------------------------------------
    # Split into chunks
    # --------------------------------------------------------

    chunks = chunk_text(
        text=text,
        chunk_size=settings.CHUNK_SIZE,
        chunk_overlap=settings.CHUNK_OVERLAP,
    )

    if not chunks:
        raise HTTPException(
            status_code=422,
            detail=(
                "The document could not be divided "
                "into usable chunks."
            ),
        )

    # --------------------------------------------------------
    # Prepare MongoDB documents
    # --------------------------------------------------------

    chunk_documents: list[
        dict[str, Any]
    ] = []

    for index, content in enumerate(
        chunks
    ):
        chunk_documents.append(
            {
                "user_id": user_id,
                "source_id": source_id,
                "chunk_index": index,
                "content": content,
                "embedding": None,
                "metadata": {
                    "chunk_size": len(
                        content
                    ),
                    "source_filename": source[
                        "filename"
                    ],
                },
            }
        )

    # --------------------------------------------------------
    # Generate embeddings for all chunks
    # --------------------------------------------------------

    chunk_contents = [
        doc["content"]
        for doc in chunk_documents
    ]
    embeddings = get_embeddings_batch(
        chunk_contents
    )

    for doc, embedding in zip(
        chunk_documents,
        embeddings,
    ):
        doc["embedding"] = embedding

    # --------------------------------------------------------
    # Store chunks
    # --------------------------------------------------------

    return create_knowledge_chunks(
        chunk_documents
    )


# ============================================================
# Get Existing Chunks
# ============================================================

def get_source_chunks(
    source_id: str,
    user_id: str,
) -> list[dict[str, Any]]:
    """
    Return all chunks for a user's knowledge source.

    Chunks are returned in their original document order.
    """

    from database import to_object_id

    source_object_id = to_object_id(
        source_id
    )

    user_object_id = to_object_id(
        user_id
    )

    documents = knowledge_chunks.find(
        {
            "source_id": source_object_id,
            "user_id": user_object_id,
        }
    ).sort(
        "chunk_index",
        1,
    )

    results: list[
        dict[str, Any]
    ] = []

    for document in documents:
        results.append(
            {
                "id": str(
                    document["_id"]
                ),
                "source_id": str(
                    document["source_id"]
                ),
                "user_id": str(
                    document["user_id"]
                ),
                "chunk_index": document[
                    "chunk_index"
                ],
                "content": document[
                    "content"
                ],
                "embedding": document.get(
                    "embedding"
                ),
                "metadata": document.get(
                    "metadata",
                    {},
                ),
                "created_at": document[
                    "created_at"
                ].isoformat(),
            }
        )

    return results


# ============================================================
# Chunk Knowledge Graph
# ============================================================

def get_user_chunk_graph(
    user_id: str,
    similarity_threshold: float = 0.86,
) -> dict[str, list[dict[str, Any]]]:
    """
    Build a graph connecting a user's document chunks.

    Nodes are chunks. Edges connect consecutive chunks within the
    same source ("sequence") and semantically similar chunks,
    including across different files ("similarity").
    """

    from database import to_object_id

    documents = list(
        knowledge_chunks.find(
            {"user_id": to_object_id(user_id)},
            {
                "source_id": 1,
                "chunk_index": 1,
                "content": 1,
                "embedding": 1,
                "metadata": 1,
            },
        ).limit(300)
    )

    nodes = [
        {
            "id": str(doc["_id"]),
            "source_id": str(doc["source_id"]),
            "filename": doc.get("metadata", {}).get(
                "source_filename", "file"
            ),
            "chunk_index": doc["chunk_index"],
            "preview": doc["content"][:80],
        }
        for doc in documents
    ]

    edges: list[dict[str, Any]] = []

    by_source: dict[str, list[dict[str, Any]]] = {}
    for doc in documents:
        by_source.setdefault(str(doc["source_id"]), []).append(doc)

    for chunks_in_source in by_source.values():
        ordered = sorted(chunks_in_source, key=lambda c: c["chunk_index"])
        for a, b in zip(ordered, ordered[1:]):
            edges.append(
                {"source": str(a["_id"]), "target": str(b["_id"]), "type": "sequence"}
            )

    embedded = [doc for doc in documents if doc.get("embedding")]
    best: dict[int, tuple[float, int]] = {}
    for i in range(len(embedded)):
        for j in range(i + 1, len(embedded)):
            similarity = _cosine_similarity(
                embedded[i]["embedding"], embedded[j]["embedding"]
            )
            if similarity < similarity_threshold:
                continue
            if similarity > best.get(i, (0.0, -1))[0]:
                best[i] = (similarity, j)
            if similarity > best.get(j, (0.0, -1))[0]:
                best[j] = (similarity, i)

    seen: set[tuple[int, int]] = set()
    for i, (similarity, j) in best.items():
        pair = (min(i, j), max(i, j))
        if pair in seen:
            continue
        seen.add(pair)
        edges.append(
            {
                "source": str(embedded[i]["_id"]),
                "target": str(embedded[j]["_id"]),
                "type": "similarity",
            }
        )

    return {"nodes": nodes, "edges": edges}


def _cosine_similarity(a: list[float], b: list[float]) -> float:
    dot = sum(x * y for x, y in zip(a, b))
    norm_a = sum(x * x for x in a) ** 0.5
    norm_b = sum(y * y for y in b) ** 0.5
    if norm_a == 0 or norm_b == 0:
        return 0.0
    return dot / (norm_a * norm_b)